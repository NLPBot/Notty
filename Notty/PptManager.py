from pptx import Presentation
from ProcessManager import ProcessManager
from SlideManager import SlideManager
from KnowledgeGraph import K_Graph
from pptx.util import Inches, Pt
import pickle


# This class manages the powerpoint slides 
# Functionalities: 
# 1. It adds slides to file 
# 2. Takes inputs to make slides

class PptManager(object):
	
	"""docstring for PptManager"""
	def __init__(self):
		self.prs = Presentation()
		self.k_graph = K_Graph()
		self.k_graph.save()


	def addSlide(self,input):
		# Create layout
		bullet_slide_layout = self.prs.slide_layouts[1]
		slide = self.prs.slides.add_slide(bullet_slide_layout) 
		# Create shape
		shapes = slide.shapes
		title_shape = shapes.title
		body_shape = shapes.placeholders[1]
		# Add slide elements
		title_shape.text = input
		tf = body_shape.text_frame
		# Create points
		processor = ProcessManager(input) # Process input
		processor.setSlideElements()
		elements = processor.getSlideElements() # elements are { (l1,c1), (l2,c2), ...   }
		# Traverse through elements to add slide components ... 
		# point = (level,content)
		for point in elements:
			p = tf.add_paragraph()
			p.font.size = Pt(20)
			p.level = point[0]
			p.text = point[1]

	def savePPT(self):
		self.prs.save('test.pptx')

import sys
# Testing...
if __name__ == '__main__':
	PPT = PptManager()
	PPT.addSlide(str(sys.argv[1]))
	#PPT.addSlide(str(sys.argv[2]))
	PPT.savePPT()







